import os
from pdfminer.high_level import extract_text
import docx
import pandas as pd
import chardet
from pptx import Presentation
from openpyxl import load_workbook

def read_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == ".txt":
            return read_txt(file_path)
        elif ext == ".csv":
            return read_csv(file_path)
        elif ext == ".pdf":
            return read_pdf(file_path)
        elif ext == ".docx":
            return read_docx(file_path)
        elif ext == ".pptx":
            return read_pptx(file_path)
        elif ext == ".xlsx":
            return read_xlsx(file_path)
        else:
            return f"Unsupported file type: {ext}"
    except Exception as e:
        return f"Error: {str(e)}"

def read_txt(path):
    with open(path, 'rb') as f:
        raw = f.read()
        encoding = chardet.detect(raw)['encoding']
    with open(path, 'r', encoding=encoding or "utf-8") as f:
        return f.read()

def read_csv(path):
    df = pd.read_csv(path)
    return df.to_string(index=False)

def read_pdf(path):
    return extract_text(path)

def read_docx(path):
    doc = docx.Document(path)
    return "\n".join([p.text for p in doc.paragraphs])

def read_pptx(path):
    prs = Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def read_xlsx(path):
    wb = load_workbook(path, data_only=True)
    output = ""
    for sheet in wb.worksheets:
        output += f"Sheet: {sheet.title}\n"
        for row in sheet.iter_rows(values_only=True):
            output += "\t".join([str(cell or "") for cell in row]) + "\n"
    return output
